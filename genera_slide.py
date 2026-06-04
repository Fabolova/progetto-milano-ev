"""
Genera una presentazione PowerPoint sul motore stocastico del Digital Twin Milano EV.
Eseguire con: .venv/bin/python genera_slide.py
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Palette colori ──
BG_DARK   = RGBColor(0x0F, 0x17, 0x2A)
ACCENT    = RGBColor(0x00, 0xD4, 0xAA)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
GRAY      = RGBColor(0xAA, 0xBB, 0xCC)
ORANGE    = RGBColor(0xFF, 0xA7, 0x26)
BLUE_L    = RGBColor(0x38, 0xBD, 0xF8)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

# ── Helper functions ──
def _add_bg(slide):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG_DARK

def _tb(slide, left, top, w, h, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.text = text; p.alignment = align
    r = p.runs[0]; r.font.size = Pt(size); r.font.color.rgb = color; r.font.bold = bold
    r.font.name = "Calibri"
    return tb

def _bullet_slide(slide, left, top, w, h, items, size=16, color=WHITE):
    tb = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    for i, txt in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = txt; p.space_after = Pt(8); p.alignment = PP_ALIGN.LEFT
        p.level = 1 if txt.startswith("   ") else 0
        if p.runs:
            r = p.runs[0]; r.font.size = Pt(size); r.font.color.rgb = color; r.font.name = "Calibri"
    return tb

def _accent_bar(slide, top):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(top), Inches(0.08), Inches(0.5))
    shape.fill.solid(); shape.fill.fore_color.rgb = ACCENT; shape.line.fill.background()

# ═══════════════════════════════════════════════════════════
# SLIDE 1 – Titolo
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 1.5, 1.5, 10, 1.2, "Milano EV Digital Twin", 44, ACCENT, True, PP_ALIGN.CENTER)
_tb(sl, 1.5, 2.8, 10, 0.8, "Motore Stocastico & Architettura MongoDB", 26, WHITE, False, PP_ALIGN.CENTER)
_tb(sl, 1.5, 4.0, 10, 0.6, "Big Data Processing & Data Engineering — A.A. 2024/25", 16, GRAY, False, PP_ALIGN.CENTER)
_tb(sl, 1.5, 5.0, 10, 0.6, "Università degli Studi di Milano-Bicocca", 14, GRAY, False, PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════
# SLIDE 2 – Agenda
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 6, 0.8, "Agenda", 32, ACCENT, True)
_bullet_slide(sl, 1.0, 1.6, 10, 5, [
    "1.  Perché MongoDB? — Vantaggi architetturali per il progetto",
    "2.  Il Motore Stocastico — Panoramica del modello probabilistico",
    "3.  Probabilità Base per Quartiere (NIL)",
    "4.  Moltiplicatore Orario — Fasce di traffico",
    "5.  Moltiplicatore Stagionale — Effetto meteo",
    "6.  Formula Composita — P(ricarica) = P_base × M_ora × M_meteo",
    "7.  Selezione Veicolo — Campionamento pesato per market share",
    "8.  Sessione di Ricarica — Calcolo energia, durata e ricavo",
    "9.  Stato IDLE — Sosta abusiva post-ricarica",
    "10. Pipeline Completa — Dal dato grezzo al business insight",
], 17, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 3 – Perché MongoDB (preambolo)
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Perché MongoDB?", 32, ACCENT, True)
_accent_bar(sl, 1.5)
_tb(sl, 1.2, 1.5, 5, 0.5, "Schema Flessibile", 20, ORANGE, True)
_bullet_slide(sl, 1.2, 2.1, 5.2, 1.5, [
    "Documenti JSON annidati (battery_specs, charging_specs)",
    "Evoluzione dello schema senza migrazioni DDL",
    "Ogni collezione ha la struttura adatta al suo dominio",
], 14, WHITE)

_accent_bar(sl, 3.5)
_tb(sl, 1.2, 3.5, 5, 0.5, "Time Series Collections", 20, ORANGE, True)
_bullet_slide(sl, 1.2, 4.1, 5.2, 1.5, [
    "Compressione nativa per dati temporali (charging_events)",
    "Bucketing automatico per query su intervalli orari",
    "Riduzione storage fino al 90% vs collezioni standard",
], 14, WHITE)

_accent_bar(sl, 5.5)
_tb(sl, 1.2, 5.5, 5, 0.5, "Atlas Cloud & Scalabilità", 20, ORANGE, True)
_bullet_slide(sl, 1.2, 6.1, 5.2, 1.2, [
    "Zero-ops: backup, replica set e monitoring integrati",
    "insert_many() per batch da 5 000 documenti alla volta",
], 14, WHITE)

# Colonna destra – vantaggi business
_tb(sl, 7.2, 1.5, 5, 0.5, "Impatto Business", 20, BLUE_L, True)
_bullet_slide(sl, 7.2, 2.1, 5, 5, [
    "✓  Time-to-market ridotto: nessuna fase di schema design",
    "✓  Costi infrastrutturali contenuti (free tier Atlas M0)",
    "✓  Indici 2dsphere per query geospaziali real-time",
    "✓  Upsert idempotenti: pipeline rieseguibili senza duplicati",
    "✓  Aggregation Framework per analytics server-side",
    "✓  Integrazione naturale con Python (PyMongo + dicts)",
], 15, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 4 – Overview Motore Stocastico
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Il Motore Stocastico — Panoramica", 32, ACCENT, True)
_bullet_slide(sl, 1.0, 1.6, 11, 5.5, [
    "Obiettivo: simulare un anno intero (2025) di sessioni di ricarica EV su Milano",
    "",
    "Approccio Discrete-Event Simulation con tick fisso di 15 minuti",
    "   → 35 040 tick simulati × N colonnine (centinaia di prese fisiche)",
    "",
    "Tre livelli di probabilità composti in un unico valore:",
    "   P(ricarica) = P_base(quartiere) × M(fascia oraria) × M(stagione)",
    "",
    "Ad ogni tick, per ogni colonnina in stato AVAILABLE:",
    "   → Estrazione uniforme U(0,1): se U < P(ricarica) → avvia sessione",
    "",
    "Tutti i parametri economici e i modelli auto letti da MongoDB (data-driven)",
], 17, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 5 – Probabilità Base per NIL
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Probabilità Base — Quartieri (NIL)", 32, ACCENT, True)
_tb(sl, 1.0, 1.5, 11, 0.6, "Ogni colonnina ha una probabilità base assegnata in base al suo NIL (Nucleo di Identità Locale):", 17, GRAY)

# Tabella
rows, cols = 4, 3
tbl_shape = sl.shapes.add_table(rows, cols, Inches(1.5), Inches(2.4), Inches(10), Inches(2.2))
tbl = tbl_shape.table
headers = ["Fascia Quartiere", "Codice NIL", "P_base (per tick 15 min)"]
data = [
    ["Centro / Alta densità",  "NIL ≤ 9",       "12%  (0.12)"],
    ["Semi-centro",            "10 ≤ NIL ≤ 45",  "8%   (0.08)"],
    ["Periferia",              "NIL > 45",        "4%   (0.04)"],
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = BG_DARK; p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri+1, ci); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3C)

_tb(sl, 1.0, 5.0, 11, 1.5, "💡 Razionale: le colonnine in centro (Duomo, Brera) hanno traffico pedonale e veicolare\n"
    "superiore, quindi una maggiore probabilità di essere utilizzate in ogni finestra di 15 minuti.", 15, GRAY)

# ═══════════════════════════════════════════════════════════
# SLIDE 6 – Moltiplicatore Orario
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Moltiplicatore Orario", 32, ACCENT, True)
_tb(sl, 1.0, 1.4, 11, 0.6, "Modella il comportamento umano: ricariche concentrate nelle ore di punta, quasi assenti di notte.", 17, GRAY)

rows, cols = 5, 3
tbl_shape = sl.shapes.add_table(rows, cols, Inches(1.5), Inches(2.2), Inches(10), Inches(2.5))
tbl = tbl_shape.table
headers = ["Fascia Oraria", "Ore", "Moltiplicatore"]
data = [
    ["🌙 Notte",           "00:00 – 05:59",  "×0.10  (quasi zero)"],
    ["🌅 Mattina (punta)", "07:00 – 10:00",  "×1.50"],
    ["☀️ Giorno",          "11:00 – 16:59",  "×1.00  (baseline)"],
    ["🌆 Sera (punta)",    "17:00 – 20:00",  "×1.80  (picco massimo)"],
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = BG_DARK; p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri+1, ci); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3C)

_tb(sl, 1.0, 5.2, 11, 1.5, "📊 Esempio: colonnina in centro (P_base=0.12) alle 18:00 → P = 0.12 × 1.80 = 21.6% per tick", 16, ORANGE)

# ═══════════════════════════════════════════════════════════
# SLIDE 7 – Moltiplicatore Stagionale
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Moltiplicatore Stagionale / Meteo", 32, ACCENT, True)
_tb(sl, 1.0, 1.4, 11, 0.6, "In inverno le batterie soffrono il freddo → più soste in ricarica. In estate il turismo aumenta la domanda.", 17, GRAY)

rows, cols = 4, 3
tbl_shape = sl.shapes.add_table(rows, cols, Inches(2), Inches(2.2), Inches(9), Inches(2))
tbl = tbl_shape.table
headers = ["Stagione", "Mesi", "Moltiplicatore"]
data = [
    ["❄️ Inverno",       "Dic, Gen, Feb",   "×1.25"],
    ["🌸 Primavera/Autunno", "Mar–Mag, Set–Nov", "×1.00"],
    ["☀️ Estate",        "Giu, Lug, Ago",   "×1.10"],
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(15); p.font.bold = True; p.font.color.rgb = BG_DARK; p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri+1, ci); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(14); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3C)

_tb(sl, 1.0, 4.8, 11, 1.5, "💡 La combinazione dei moltiplicatori crea pattern realistici:\n"
    "una sera di gennaio in centro raggiunge P = 0.12 × 1.80 × 1.25 = 27%", 16, ORANGE)

# ═══════════════════════════════════════════════════════════
# SLIDE 8 – Formula Composita
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Formula Composita", 32, ACCENT, True)

# Formula grande
_tb(sl, 1.5, 1.8, 10, 1, "P(ricarica)  =  P_base(NIL)  ×  M_orario(h)  ×  M_stagione(m)", 28, ORANGE, True, PP_ALIGN.CENTER)

_tb(sl, 1.0, 3.0, 11, 0.5, "Decisione per ogni tick (ogni 15 minuti):", 18, WHITE, True)
_bullet_slide(sl, 1.2, 3.6, 10, 3.5, [
    "1.  Genera un numero casuale U ~ Uniform(0, 1)",
    "2.  Se  U < P(ricarica)  →  la colonnina inizia una sessione di ricarica",
    "3.  Se  U ≥ P(ricarica)  →  la colonnina resta in stato AVAILABLE",
    "",
    "Questo è un processo di Bernoulli ad ogni tick:",
    "   → Successo (ricarica) con probabilità p",
    "   → Insuccesso (resta libera) con probabilità (1 − p)",
    "",
    "Il tempo di attesa tra due ricariche segue una distribuzione Geometrica",
], 16, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 9 – Selezione Veicolo
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Selezione Veicolo — Campionamento Pesato", 32, ACCENT, True)
_tb(sl, 1.0, 1.4, 11, 0.6, "Il modello dell'auto viene estratto con random.choices() usando pesi di market share:", 17, GRAY)

rows, cols = 6, 3
tbl_shape = sl.shapes.add_table(rows, cols, Inches(1.5), Inches(2.2), Inches(10), Inches(3))
tbl = tbl_shape.table
headers = ["Modello", "Tipo", "Peso Market Share"]
data = [
    ["Tesla Model Y LR",    "BEV",  "15%"],
    ["Fiat 500e 42kWh",     "BEV",  "15%"],
    ["Tesla Model 3 RWD",   "BEV",  "10%"],
    ["Smart EQ fortwo",     "BEV",  "10%  (solo AC!)"],
    ["Jeep Compass 4xe",    "PHEV", "10%  (no DC, lento)"],
]
for ci, h in enumerate(headers):
    cell = tbl.cell(0, ci); cell.text = h
    for p in cell.text_frame.paragraphs:
        p.font.size = Pt(14); p.font.bold = True; p.font.color.rgb = BG_DARK; p.alignment = PP_ALIGN.CENTER
    cell.fill.solid(); cell.fill.fore_color.rgb = ACCENT
for ri, row in enumerate(data):
    for ci, val in enumerate(row):
        cell = tbl.cell(ri+1, ci); cell.text = val
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(13); p.font.color.rgb = WHITE; p.alignment = PP_ALIGN.CENTER
        cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x1A, 0x25, 0x3C)

_tb(sl, 1.0, 5.5, 11, 1.5, "⚠️ Vincolo di compatibilità: se la colonnina è DC e il veicolo ha max_dc_kw = 0, la sessione viene scartata.", 15, ORANGE)

# ═══════════════════════════════════════════════════════════
# SLIDE 10 – Sessione di Ricarica
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Sessione di Ricarica — Calcolo Durata e Ricavo", 32, ACCENT, True)
_bullet_slide(sl, 1.0, 1.6, 11, 5.5, [
    "1.  SOC iniziale estratto casualmente:  SOC₀ ~ Uniform(10%, 50%)",
    "",
    "2.  Energia necessaria:  kWh = Capacità_batteria × (80% − SOC₀) / 100",
    "     → Si ricarica sempre fino all'80% (best practice per longevità batteria)",
    "",
    "3.  Potenza effettiva:  P_eff = min(P_colonnina, P_max_veicolo)",
    "     → Il collo di bottiglia può essere la colonnina O il veicolo",
    "",
    "4.  Durata:  T = kWh / P_eff   (convertita in minuti)",
    "",
    "5.  Ricavo per tick:  R_tick = P_eff × (15/60) × Tariffa_€/kWh",
    "     → AC: €0.65/kWh  |  DC: €0.90/kWh  (da simulation_parameters)",
    "",
    "6.  Il ricavo totale si accumula tick dopo tick fino al completamento",
], 16, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 11 – Stato IDLE
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Stato IDLE — Sosta Abusiva Post-Ricarica", 32, ACCENT, True)
_bullet_slide(sl, 1.0, 1.6, 11, 5, [
    "Al termine di ogni sessione di ricarica:",
    "",
    "   → Con probabilità 20% il veicolo resta parcheggiato (stato IDLE)",
    "   → Con probabilità 80% il veicolo se ne va (stato → AVAILABLE)",
    "",
    "Se il veicolo resta in sosta abusiva:",
    "   • Durata IDLE ~ Uniform(15, 60) minuti",
    "   • Si applica una penale di €0.10/minuto (idle_fee)",
    "   • Il ricavo della penale viene tracciato come evento separato",
    "",
    "💡 Questo modella un fenomeno reale e misurabile:",
    "   la sosta abusiva è la causa #1 di sotto-utilizzo delle colonnine",
    "   e rappresenta un'opportunità di revenue aggiuntiva per l'operatore",
], 16, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 12 – Pipeline Completa
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 0.8, 0.5, 10, 0.8, "Pipeline End-to-End", 32, ACCENT, True)
_bullet_slide(sl, 1.0, 1.5, 5.5, 5.5, [
    "① Censimento Geo (GeoJSON → MongoDB)",
    "   Colonnine + coordinate + NIL + titolare",
    "",
    "② Enrichment Demografico",
    "   Densità abitanti/km² per quartiere",
    "",
    "③ Parametri Economici",
    "   Tariffe B2C, costi B2B (PUN + oneri)",
    "",
    "④ Catalogo EV (10 modelli con market share)",
    "",
    "⑤ Motore Stocastico (07_connected_simulation)",
    "   35 040 tick × N prese → charging_sessions_log",
], 15, WHITE)
_bullet_slide(sl, 7.0, 1.5, 5.5, 5.5, [
    "⑥ Dati Meteo Storici (Open-Meteo API)",
    "   Pioggia oraria 2023–2026 → CSV",
    "",
    "⑦ Grafo Neo4j",
    "   Colonnine → Quartieri → Eventi → Veicoli",
    "   POI, Musei, Supermercati di prossimità",
    "",
    "⑧ Import Transazioni (batch UNWIND)",
    "   MongoDB → Neo4j con prestazioni ottimizzate",
    "",
    "⑨ Analytics & Digital Twin",
    "   Query geospaziali, tassi di occupazione,",
    "   redditività per quartiere e fascia oraria",
], 15, WHITE)

# ═══════════════════════════════════════════════════════════
# SLIDE 13 – Chiusura
# ═══════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
_add_bg(sl)
_tb(sl, 1.5, 2.0, 10, 1, "Grazie per l'attenzione", 40, ACCENT, True, PP_ALIGN.CENTER)
_tb(sl, 1.5, 3.5, 10, 0.8, "Milano EV Digital Twin — Big Data Processing & Data Engineering", 20, WHITE, False, PP_ALIGN.CENTER)
_tb(sl, 1.5, 4.5, 10, 0.6, "Domande?", 28, ORANGE, True, PP_ALIGN.CENTER)

# ── Salvataggio ──
OUTPUT = "presentazione_motore_stocastico.pptx"
prs.save(OUTPUT)
print(f"✅ Presentazione salvata in '{OUTPUT}' ({len(prs.slides)} slide)")
